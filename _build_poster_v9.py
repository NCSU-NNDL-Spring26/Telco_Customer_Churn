"""poster_v9 builder.

Default: opens existing poster_v9.pptx, replaces **section A (Motivation)** only
(__V9A__-tagged shapes + anything whose center lies in column A row 1), and keeps
your manual edits in B–G and the header strip.

Full regeneration from poster_v3.pptx (wipes the slide except KEEP shapes):

  set POSTER_V9_FULL_REBUILD=1

Cost curve layout: ~$370K tiny flush right; counterfactual strip; F/G sizing."""
import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE_PPTX = "poster_v3.pptx"
DST = "poster_v9.pptx"
FIG = "figures"

EMU_PER_INCH = 914400
AUTO_A_PREFIX = "__V9A__"

# --- palette ---------------------------------------------------------------
NAVY  = RGBColor(0x2C, 0x3E, 0x50)
TEAL  = RGBColor(0x1F, 0x6F, 0x7A)
WINE  = RGBColor(0x7B, 0x2D, 0x40)
DARK  = RGBColor(0x21, 0x25, 0x29)
GRAY  = RGBColor(0x5F, 0x6A, 0x76)
LGRAY = RGBColor(0xE8, 0xEB, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

A_BG = RGBColor(0xEC,0xEE,0xF2); A_HDR = RGBColor(0xC8,0xD2,0xDD)
B_BG = RGBColor(0xE2,0xF1,0xF4); B_HDR = RGBColor(0xB8,0xDC,0xE2)
C_BG = RGBColor(0xE8,0xF0,0xE0); C_HDR = RGBColor(0xC4,0xD5,0xB5)
D_BG = RGBColor(0xF6,0xEE,0xD8); D_HDR = RGBColor(0xE6,0xD3,0x9B)
E_BG = RGBColor(0xE5,0xEB,0xF1); E_HDR = RGBColor(0xBD,0xCD,0xDD)
F_BG = RGBColor(0xED,0xE7,0xF4); F_HDR = RGBColor(0xC9,0xBA,0xE0)
G_BG = RGBColor(0xDD,0xEE,0xEC); G_HDR = RGBColor(0xA8,0xD0,0xCA)

TBL_HDR_BG = NAVY
TBL_WIN_BG = RGBColor(0xDC,0xED,0xC8)
TBL_ALT_BG = RGBColor(0xF8,0xF9,0xFA)
TBL_RW_BG  = WHITE

# --- helpers ---------------------------------------------------------------
def E(x): return Inches(x)
def remove_shape(sh): sh._element.getparent().remove(sh._element)

def add_text(slide, x, y, w, h, lines, size, *, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None, font="Calibri"):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align
        if line_spacing is not None: p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.italic = italic; r.font.name = font
            r.font.color.rgb = color
    return tb

def add_panel(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def add_header_bar(slide, x, y, w, h, fill_rgb, text, size=38, color=DARK):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = color; r.font.name = "Calibri"
    return sh

def add_image_fit(slide, x, y, max_w, max_h, path):
    iw, ih = Image.open(path).size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h, w = max_h, max_h * ratio
    cx = x + (max_w - w) / 2; cy = y + (max_h - h) / 2
    return slide.shapes.add_picture(path, E(cx), E(cy), E(w), E(h))

# --- layout ----------------------------------------------------------------
SW, SH = 48.0, 36.0
M = 0.40
GAP = 0.40

A_W = 14.20
B_W = 18.00
C_W = SW - 2*M - 2*GAP - A_W - B_W   # 14.20

X_A = M
X_B = X_A + A_W + GAP
X_C = X_B + B_W + GAP

ROW1_Y, ROW1_H = 5.70, 10.80
ROW2_Y, ROW2_H = 16.70, 11.80
ROW3_Y, ROW3_H = 28.70,  5.60
REF_Y,  REF_H  = 34.50,  1.20

COLW2 = (SW - 2*M - GAP) / 2
X_L = M
X_R = M + COLW2 + GAP

HDR_H = 1.20
PAD_X = 0.40

KEEP = {"Picture 10", "TextBox 2", "TextBox 3", "TextBox 1",
        "Title 1", "Subtitle 2"}


def _shape_center_inch(sh):
    return (
        (sh.left + sh.width / 2) / EMU_PER_INCH,
        (sh.top + sh.height / 2) / EMU_PER_INCH,
    )


def remove_motivation_section(slide, x_a, y_r1, w_a, h_r1):
    """Strip column A row 1 so Motivation can be rebuilt without touching B–G."""
    victims = []
    for sh in list(slide.shapes):
        if sh.name in KEEP:
            continue
        if str(sh.name).startswith(AUTO_A_PREFIX):
            victims.append(sh)
            continue
        cx, cy = _shape_center_inch(sh)
        if x_a <= cx <= x_a + w_a and y_r1 <= cy <= y_r1 + h_r1:
            victims.append(sh)
    for sh in victims:
        remove_shape(sh)


def patch_nndl_header(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and ("Poster 60" in sh.text_frame.text
                                  or "Poster ###" in sh.text_frame.text):
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = "NNDL.39"
            break


_full_rebuild = os.environ.get("POSTER_V9_FULL_REBUILD", "").strip().lower() in (
    "1", "true", "yes", "on",
)
merge_a_only = (not _full_rebuild) and os.path.isfile(DST)

if merge_a_only:
    prs = Presentation(DST)
    slide = prs.slides[0]
    remove_motivation_section(slide, X_A, ROW1_Y, A_W, ROW1_H)
else:
    if not os.path.isfile(BASE_PPTX):
        raise FileNotFoundError(
            f"{BASE_PPTX} not found; required when poster_v9.pptx is missing "
            "or POSTER_V9_FULL_REBUILD=1."
        )
    shutil.copy(BASE_PPTX, DST)
    prs = Presentation(DST)
    slide = prs.slides[0]
    for sh in [s for s in slide.shapes if s.name not in KEEP]:
        remove_shape(sh)
    patch_nndl_header(slide)

# =================================================================
# A.  MOTIVATION  (always rebuilt; shapes tagged for next merge)
# =================================================================
_sh = add_panel(slide, X_A, ROW1_Y, A_W, ROW1_H, A_BG); _sh.name = AUTO_A_PREFIX + "panel"
_sh = add_header_bar(slide, X_A, ROW1_Y, A_W, HDR_H, A_HDR, "A.  Motivation", 38)
_sh.name = AUTO_A_PREFIX + "hdr"

cx = X_A + PAD_X; cy = ROW1_Y + HDR_H + 0.15; cw = A_W - 2*PAD_X

_sh = add_text(slide, cx, cy, cw, 2.45, [
    "Customer churn -- when a paying subscriber leaves a service -- is one of the most expensive problems in telecom.",
    "Acquiring a new customer costs 5x-7x more than retaining one, so identifying customers about to leave is directly tied to revenue.",
    "Churn is rarely a surprise overnight: usage slips, support tickets spike, bills become unpredictable -- weak signals appear weeks or months before cancel.",
    "Carriers already hold the data to spot those patterns; the gap is reliable risk scores and thresholds aligned with who you can afford to save.",
], size=17, color=DARK, line_spacing=1.12)
_sh.name = AUTO_A_PREFIX + "hook"
cy += 2.58

_sh = add_text(slide, cx, cy, cw, 1.0, [
    "For a 1M-subscriber carrier with 25% churn, a 3% reduction = ~7,500 retained customers/year, worth millions in lifetime value.",
], size=16, color=WINE, italic=True, line_spacing=1.15)
_sh.name = AUTO_A_PREFIX + "stakes"
cy += 1.10

_sh = add_text(slide, cx, cy, cw, 0.50, ["Why standard ML misses the mark"],
               size=22, bold=True, color=NAVY)
_sh.name = AUTO_A_PREFIX + "why_hdr"
cy += 0.55

chips = [
    ("Imbalanced labels",
     "Only 26.6% of customers churn. A trivial 'predict no-churn' model gets 73% accuracy and zero business value. Accuracy is the wrong metric."),
    ("Asymmetric cost",
     "Missing a churner ($1,500 lost LTV) is ~30x worse than a false alarm ($50 retention call). Optimization must reflect this."),
    ("Threshold blindness",
     "Default 0.5 cutoffs ignore both imbalance and cost. A 'good' model with the wrong threshold is still a poor business tool."),
]
chip_h = 1.45
for i, (label, body) in enumerate(chips):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                E(cx), E(cy), E(cw), E(chip_h))
    sh.name = f"{AUTO_A_PREFIX}chip{i}"
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = NAVY; sh.line.width = Pt(1.0)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(180000); tf.margin_right = Emu(150000)
    tf.margin_top = Emu(80000); tf.margin_bottom = Emu(80000)
    p = tf.paragraphs[0]; p.text = label
    for r in p.runs:
        r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = NAVY; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.text = body; p2.line_spacing = 1.15
    for r in p2.runs:
        r.font.size = Pt(15); r.font.color.rgb = DARK
        r.font.name = "Calibri"
    cy += chip_h + 0.15

# =================================================================
if not merge_a_only:
    # B.  DATA & EDA
    # =================================================================
    add_panel(slide, X_B, ROW1_Y, B_W, ROW1_H, B_BG)
    add_header_bar(slide, X_B, ROW1_Y, B_W, HDR_H, B_HDR, "B.  Data & EDA", 38)

    cx = X_B + PAD_X; cy = ROW1_Y + HDR_H + 0.15; cw = B_W - 2*PAD_X

    add_text(slide, cx, cy, cw, 0.50, ["Dataset:  IBM Telco Customer Churn  (Kaggle)"],
             size=22, bold=True, color=NAVY); cy += 0.55
    add_text(slide, cx, cy, cw, 1.4, [
        "7,032 customers x 19 features after cleaning (11 dropped for invalid TotalCharges).",
        "Demographics, tenure, contract, services (phone/internet/security/streaming), monthly + total charges.   Target: Churn (Yes/No).",
    ], size=17, color=DARK, line_spacing=1.15)
    cy += 1.55

    # Class balance + churn-by-contract side by side
    fig_h = 4.10
    left_w = cw * 0.42
    right_w = cw - left_w - 0.20
    add_image_fit(slide, cx, cy, left_w, fig_h, f"{FIG}/eda_01_class_balance.png")
    add_image_fit(slide, cx + left_w + 0.20, cy, right_w, fig_h,
                  f"{FIG}/eda_02_churn_by_contract.png")
    cy += fig_h + 0.10

    add_text(slide, cx, cy, cw, 0.50, ["What the data tells us"],
             size=22, bold=True, color=NAVY); cy += 0.55
    add_text(slide, cx, cy, cw, 2.5, [
        "Imbalanced class:  5,163 stayed (73.4%)  vs  1,869 churned (26.6%)  --  ratio 1 : 2.76.",
        "Contract dominates:  month-to-month 43%  vs  2-year 3%  --  a 14x gap.",
        "Tenure:  churners cluster heavily in the first 12 months.",
        "Service:  fiber 42%  vs  DSL 19%;   electronic-check 45%  vs  auto-pay <20%.",
    ], size=16, color=DARK, line_spacing=1.20)

    # =================================================================
    # C.  METHODOLOGY
    # =================================================================
    add_panel(slide, X_C, ROW1_Y, C_W, ROW1_H, C_BG)
    add_header_bar(slide, X_C, ROW1_Y, C_W, HDR_H, C_HDR, "C.  Methodology", 38)

    cx = X_C + PAD_X; cy = ROW1_Y + HDR_H + 0.10; cw = C_W - 2*PAD_X

    add_text(slide, cx, cy, cw, 0.95, [
        "One disciplined pipeline -- so differences come from architecture, not preprocessing.",
    ], size=17, color=DARK, italic=True, line_spacing=1.15)
    cy += 0.95

    def flow_box(x, y, w, h, text, fill, size=14, color=DARK):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    E(x), E(y), E(w), E(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = NAVY; sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        tf = sh.text_frame; tf.word_wrap = True
        tf.margin_top = Emu(30000); tf.margin_bottom = Emu(30000)
        tf.margin_left = Emu(60000); tf.margin_right = Emu(60000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = color; r.font.name = "Calibri"

    def flow_arrow(x_center, y_top, h=0.18):
        sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    E(x_center - 0.14), E(y_top), E(0.28), E(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = NAVY
        sh.line.fill.background()

    bw = cw; fy = cy
    flow_box(cx, fy, bw, 0.50, "Raw CSV  (7,032 customers)", WHITE, 16); fy += 0.50
    flow_arrow(cx + bw/2, fy); fy += 0.18
    flow_box(cx, fy, bw, 0.55, "Clean  +  one-hot encode", WHITE, 14); fy += 0.55
    flow_arrow(cx + bw/2, fy); fy += 0.18
    flow_box(cx, fy, bw, 0.65,
             "Stratified split  70 / 15 / 15\nFit StandardScaler on TRAIN only",
             RGBColor(0xFF,0xF8,0xDC), 13); fy += 0.65
    flow_arrow(cx + bw/2, fy); fy += 0.18
    mb_w = (bw - 0.30) / 4
    flow_box(cx + 0*(mb_w+0.10), fy, mb_w, 0.55, "Logistic\nReg.",  WHITE, 12)
    flow_box(cx + 1*(mb_w+0.10), fy, mb_w, 0.55, "Custom\nMLP",     WHITE, 12)
    flow_box(cx + 2*(mb_w+0.10), fy, mb_w, 0.55, "Cat\nBoost",      WHITE, 12)
    flow_box(cx + 3*(mb_w+0.10), fy, mb_w, 0.55, "XG\nBoost",       WHITE, 12)
    fy += 0.55
    flow_arrow(cx + bw/2, fy); fy += 0.18
    flow_box(cx, fy, bw, 0.50, "Tune threshold on VAL  (max F1)",
             RGBColor(0xFF,0xF8,0xDC), 14); fy += 0.50
    flow_arrow(cx + bw/2, fy); fy += 0.18
    flow_box(cx, fy, bw, 0.50, "Test metrics  +  Cost analysis",
             C_HDR, 14); fy += 0.50

    fy += 0.18
    add_text(slide, cx, fy, cw, 0.50, ["Why each step"],
             size=20, bold=True, color=NAVY); fy += 0.50
    add_text(slide, cx, fy, cw, 2.6, [
        "Stratified split:  preserves the 26.6% churn rate in every fold.",
        "Train-only scaling:  prevents test-set leakage into training.",
        "Threshold tuning:  default 0.5 is arbitrary on imbalanced data.",
        "Brier + PR-AUC:  measure calibration and minority-class ranking.",
    ], size=14, color=DARK, line_spacing=1.20)

    # =================================================================
    # D.  MODEL COMPARISON
    # =================================================================
    add_panel(slide, X_L, ROW2_Y, COLW2, ROW2_H, D_BG)
    add_header_bar(slide, X_L, ROW2_Y, COLW2, HDR_H, D_HDR,
                   "D.  Results -- Model Comparison", 40)

    cx = X_L + PAD_X; cy = ROW2_Y + HDR_H + 0.20; cw = COLW2 - 2*PAD_X

    add_text(slide, cx, cy, cw, 1.4, [
        "All four models trained on identical features and splits, scored on test with a validation-tuned threshold (max F1).",
        "Brier captures calibration:  do predicted probabilities match observed churn rates? Lower is better.",
    ], size=20, color=DARK, italic=True, line_spacing=1.15)
    cy += 1.45

    comp_rows = [
        ["Model", "F1", "ROC-AUC", "PR-AUC", "Brier", "Recall", "Precision"],
        ["CatBoost",            "0.624", "0.841", "0.642", "0.139", "0.768", "0.526"],
        ["Logistic Regression", "0.610", "0.831", "0.591", "0.176", "0.746", "0.516"],
        ["XGBoost",             "0.608", "0.838", "0.643", "0.169", "0.725", "0.523"],
        ["MLP (custom NN)",     "0.603", "0.826", "0.592", "0.177", "0.789", "0.488"],
    ]
    T_H = 2.30
    tbl = slide.shapes.add_table(len(comp_rows), 7,
                                 E(cx), E(cy), E(cw), E(T_H)).table
    for r in range(len(comp_rows)):
        for c in range(7):
            cell = tbl.cell(r, c); cell.text = comp_rows[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(19); run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True; run.font.color.rgb = WHITE
                    else: run.font.color.rgb = DARK
                    if r == 1: run.font.bold = True
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR_BG
            elif r == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_WIN_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_ALT_BG if r % 2 == 0 else TBL_RW_BG
    cy += T_H + 0.15

    add_text(slide, cx, cy, cw, 1.4, [
        "CatBoost wins on F1, ROC-AUC, and Brier -- it ranks well AND its probabilities are calibrated.",
        "MLP does not beat logistic regression -- on small one-hot tabular data, a linear boundary is most of the signal.",
    ], size=16, color=NAVY, italic=True, line_spacing=1.15)
    cy += 1.45

    fig_max_h = ROW2_Y + ROW2_H - cy - 0.50
    add_image_fit(slide, cx, cy, cw, fig_max_h, f"{FIG}/boosted_roc_pr_curves.png")
    add_text(slide, cx, ROW2_Y + ROW2_H - 0.45, cw, 0.40,
             ["ROC and Precision-Recall curves on test set."],
             size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    # =================================================================
    # E.  COST-AWARE ANALYSIS  -- HERO COST CURVE, $370K small on right
    # =================================================================
    add_panel(slide, X_R, ROW2_Y, COLW2, ROW2_H, E_BG)
    add_header_bar(slide, X_R, ROW2_Y, COLW2, HDR_H, E_HDR,
                   "E.  Cost-Aware Analysis", 40)

    cx = X_R + PAD_X; cy = ROW2_Y + HDR_H + 0.15; cw = COLW2 - 2*PAD_X

    # Narrative -- compact 1 line
    add_text(slide, cx, cy, cw, 0.75, [
        "Sweep threshold under realistic costs:  FN $1,500   FP $50.   Net Value = TP*$1500 - FP*$50.   Optimum chosen on validation.",
    ], size=16, color=DARK, italic=True, line_spacing=1.10)
    cy += 0.78

    # TOP ROW: cost optima table (wide) + tiny ~$370K flush RIGHT
    top_row_h = 1.85
    card_w = min(2.15, cw * 0.22)
    tbl_w = cw - card_w - 0.12
    card_x = cx + cw - card_w

    # Cost optima table on left
    cost_rows = [
        ["Model", "Threshold", "Net value / 1,000"],
        ["Logistic Regression", "0.11", "$371,043"],
        ["XGBoost",             "0.09", "$370,616"],
        ["CatBoost",            "0.03", "$370,142"],
        ["MLP (custom NN)",     "0.07", "$368,910"],
    ]
    ct = slide.shapes.add_table(len(cost_rows), 3,
                                E(cx), E(cy), E(tbl_w), E(top_row_h)).table
    for r in range(len(cost_rows)):
        for c in range(3):
            cell = ct.cell(r, c); cell.text = cost_rows[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(14); run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True; run.font.color.rgb = WHITE
                    else: run.font.color.rgb = DARK
                    if c == 2 and r > 0: run.font.bold = True
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_ALT_BG if r % 2 == 0 else TBL_RW_BG

    # Tiny ~$370K label flush right (does not steal space from the hero plot)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  E(card_x), E(cy), E(card_w), E(top_row_h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TEAL; card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(45000); tf.margin_right = Emu(45000)
    tf.margin_top = Emu(50000); tf.margin_bottom = Emu(50000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = "~$370K"; p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(17); r.font.bold = True
        r.font.color.rgb = TEAL; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.text = "net / 1K cust."; p2.alignment = PP_ALIGN.RIGHT
    for r in p2.runs:
        r.font.size = Pt(9); r.font.italic = True
        r.font.color.rgb = NAVY; r.font.name = "Calibri"
    p3 = tf.add_paragraph(); p3.text = "(thr.-tuned)"; p3.alignment = PP_ALIGN.RIGHT
    for r in p3.runs:
        r.font.size = Pt(8); r.font.italic = True
        r.font.color.rgb = GRAY; r.font.name = "Calibri"

    cy += top_row_h + 0.15

    # Constraint sentence
    add_text(slide, cx, cy, cw, 0.42, [
        "Top-20% capacity constraint:  CatBoost still delivers $175,687 / 1,000 at threshold 0.54.",
    ], size=14, color=WINE, italic=True, line_spacing=1.10)
    cy += 0.48

    # HERO COST CURVE -- maximize vertical space; thin strip for counterfactual
    cc_max_h = (ROW2_Y + ROW2_H) - cy - 0.98
    add_image_fit(slide, cx, cy, cw, cc_max_h, f"{FIG}/cost_curve.png")
    cy += cc_max_h + 0.10

    # Counterfactual horizontal strip (3 segments: customer -> action -> result)
    strip_h = 0.95
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   E(cx), E(cy), E(cw), E(strip_h))
    strip.fill.solid(); strip.fill.fore_color.rgb = WHITE
    strip.line.color.rgb = NAVY; strip.line.width = Pt(1.5)
    strip.shadow.inherit = False

    col_w = cw / 3
    add_text(slide, cx + 0.20, cy + 0.05, col_w - 0.30, 0.30,
             ["COUNTERFACTUAL"], size=11, bold=True, color=NAVY)
    add_text(slide, cx + 0.20, cy + 0.35, col_w - 0.30, 0.55,
             ["Customer:  predicted churn = 0.93"], size=14, color=DARK, bold=True)

    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                E(cx + col_w - 0.10), E(cy + 0.30),
                                E(0.30), E(0.30))
    a1.fill.solid(); a1.fill.fore_color.rgb = NAVY; a1.line.fill.background()

    add_text(slide, cx + col_w + 0.20, cy + 0.10, col_w - 0.30, 0.30,
             ["ACTION"], size=11, bold=True, color=NAVY)
    add_text(slide, cx + col_w + 0.20, cy + 0.40, col_w - 0.30, 0.50,
             ["Switch InternetService:  Fiber -> No internet"],
             size=13, color=DARK)

    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                E(cx + 2*col_w - 0.10), E(cy + 0.30),
                                E(0.30), E(0.30))
    a2.fill.solid(); a2.fill.fore_color.rgb = NAVY; a2.line.fill.background()

    add_text(slide, cx + 2*col_w + 0.20, cy + 0.10, col_w - 0.30, 0.30,
             ["RESULT"], size=11, bold=True, color=NAVY)
    add_text(slide, cx + 2*col_w + 0.20, cy + 0.40, col_w - 0.30, 0.50,
             ["predicted churn = 0.82  (-11 pp)"],
             size=15, bold=True, color=TEAL)

    # =================================================================
    # F.  WHAT DRIVES CHURN  (overflow fixed)
    # =================================================================
    add_panel(slide, X_L, ROW3_Y, COLW2, ROW3_H, F_BG)
    add_header_bar(slide, X_L, ROW3_Y, COLW2, HDR_H, F_HDR,
                   "F.  What Drives Churn", 36)

    cx = X_L + PAD_X; cy = ROW3_Y + HDR_H + 0.15; cw = COLW2 - 2*PAD_X

    # Image (left) + table + interpretation (right)
    fi_max_h = ROW3_Y + ROW3_H - cy - 0.20
    fi_w = cw * 0.46
    add_image_fit(slide, cx, cy, fi_w, fi_max_h, f"{FIG}/feature_importance_catboost.png")

    ft_x = cx + fi_w + 0.30; ft_w = cw - fi_w - 0.30
    feat_rows = [
        ["Top driver", "Importance"],
        ["tenure",                            "21.4"],
        ["Contract  =  Two year",             "14.4"],
        ["InternetService  =  Fiber optic",   "13.0"],
        ["TotalCharges",                      " 7.6"],
        ["MonthlyCharges",                    " 7.4"],
        ["Contract  =  One year",             " 4.8"],
        ["PaymentMethod  =  Electronic check"," 4.3"],
    ]
    ft_h = fi_max_h * 0.65
    ft = slide.shapes.add_table(len(feat_rows), 2,
                                E(ft_x), E(cy), E(ft_w), E(ft_h)).table
    for r in range(len(feat_rows)):
        for c in range(2):
            cell = ft.cell(r, c); cell.text = feat_rows[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(15); run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True; run.font.color.rgb = WHITE
                    else: run.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_ALT_BG if r % 2 == 0 else TBL_RW_BG

    # Interpretation tightened (3 short lines, fits inside panel)
    interp_y = cy + ft_h + 0.10
    interp_h = fi_max_h - ft_h - 0.10
    add_text(slide, ft_x, interp_y, ft_w, interp_h, [
        "Aligns with industry intuition:  short tenure, no long contract,",
        "fiber, electronic check  --  known retention friction points.",
        "Lever:  contract upsell campaigns target the top 3 drivers.",
    ], size=13, color=DARK, italic=True, line_spacing=1.20)

    # =================================================================
    # G.  CONCLUSIONS & NEXT STEPS  (right column simplified to fit width)
    # =================================================================
    add_panel(slide, X_R, ROW3_Y, COLW2, ROW3_H, G_BG)
    add_header_bar(slide, X_R, ROW3_Y, COLW2, HDR_H, G_HDR,
                   "G.  Conclusions & Next Steps", 36)

    cx = X_R + PAD_X; cy = ROW3_Y + HDR_H + 0.15; cw = COLW2 - 2*PAD_X

    gw_l = cw * 0.58
    gw_r = cw - gw_l - 0.25
    gx_r = cx + gw_l + 0.25

    # LEFT: Findings (4 numbered)
    add_text(slide, cx, cy, gw_l, 0.45, ["Findings"],
             size=20, bold=True, color=NAVY)
    add_text(slide, cx, cy + 0.50, gw_l, 3.7, [
        "1.  CatBoost is the strongest classifier (F1 0.624, ROC-AUC 0.841, lowest Brier 0.139).",
        "2.  Custom MLP does not beat logistic regression -- nonlinearity adds little on small one-hot data.",
        "3.  Cost-aware threshold tuning is the dominant lever -- all four models recover ~$370K per 1,000 customers; architecture changes the result by < 1%.",
        "4.  Counterfactuals turn probabilities into specific retention actions.",
    ], size=13, color=DARK, line_spacing=1.25)

    # RIGHT: 3 compact sub-blocks (shorter lines so nothing clips)
    add_text(slide, gx_r, cy, gw_r, 0.40, ["Deployment readiness"],
             size=15, bold=True, color=NAVY)
    add_text(slide, gx_r, cy + 0.40, gw_r, 1.0, [
        "Save:  model + scaler + threshold (0.31).",
        "Score monthly;  recalibrate quarterly.",
    ], size=12, color=DARK, line_spacing=1.20)

    add_text(slide, gx_r, cy + 1.55, gw_r, 0.40, ["Limitations"],
             size=15, bold=True, color=NAVY)
    add_text(slide, gx_r, cy + 1.95, gw_r, 1.0, [
        "Static snapshot;  no temporal drift.",
        "Cost ratios illustrative -- vary by segment.",
    ], size=12, color=DARK, line_spacing=1.20)

    add_text(slide, gx_r, cy + 3.10, gw_r, 0.40, ["Next steps"],
             size=15, bold=True, color=NAVY)
    add_text(slide, gx_r, cy + 3.50, gw_r, 1.0, [
        "Stacking ensemble.",
        "Cost-sensitive loss in MLP.",
        "Live A/B test of threshold.",
    ], size=12, color=DARK, line_spacing=1.20)

    # =================================================================
    # REFERENCES STRIP
    # =================================================================
    add_panel(slide, M, REF_Y, SW - 2*M, REF_H, LGRAY, line=GRAY)
    add_text(slide, M + 0.30, REF_Y + 0.05, SW - 2*M - 0.6, 0.35,
             ["References & Acknowledgements"], size=16, bold=True, color=NAVY)
    add_text(slide, M + 0.30, REF_Y + 0.40, SW - 2*M - 0.6, REF_H - 0.45, [
        "[1] Verbeke, W. et al. (2012). New insights into churn prediction in telecommunications. Eur. J. Op. Res., 218(1), 211-229.",
        "[2] Prokhorenkova, L. et al. (2018). CatBoost: unbiased boosting with categorical features. NeurIPS.",
        "[3] Chen, T. & Guestrin, C. (2016). XGBoost: a scalable tree-boosting system. KDD.",
        "Dataset:  IBM Telco Customer Churn (Kaggle: blastchar/telco-customer-churn).    Acknowledgements:  Dr. Edgar Lobaton (instructor, ECE 525, NC State University).",
    ], size=12, color=DARK, line_spacing=1.10)

prs.save(DST)
print("Wrote", DST,
      "(section A only — kept your manual edits)" if merge_a_only else "(full rebuild from poster_v3)")
