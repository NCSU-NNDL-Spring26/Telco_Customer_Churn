"""poster_v10: build on poster_v9 — keeps rows 1–2 (A–E) intact, replaces the bottom
band only with a taller row 3 and asymmetric F vs G. Section F uses a much larger
feature-importance figure.

Requires poster_v9.pptx (your edited deck). Does not modify poster_v9.

  py -3 _build_poster_v10.py

Output: poster_v10.pptx"""
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE_PPTX = "poster_v9.pptx"
DST = "poster_v10.pptx"
FIG = "figures"

EMU_PER_INCH = 914400

# --- palette (match v9) ----------------------------------------------------
NAVY  = RGBColor(0x2C, 0x3E, 0x50)
DARK  = RGBColor(0x21, 0x25, 0x29)
GRAY  = RGBColor(0x5F, 0x6A, 0x76)
LGRAY = RGBColor(0xE8, 0xEB, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_BG = RGBColor(0xED,0xE7,0xF4); F_HDR = RGBColor(0xC9,0xBA,0xE0)
G_BG = RGBColor(0xDD,0xEE,0xEC); G_HDR = RGBColor(0xA8,0xD0,0xCA)

TBL_HDR_BG = NAVY
TBL_ALT_BG = RGBColor(0xF8,0xF9,0xFA)
TBL_RW_BG  = WHITE

# --- helpers ---------------------------------------------------------------
def E(x): return Inches(x)


def remove_shape(sh):
    sh._element.getparent().remove(sh._element)


def add_text(slide, x, y, w, h, lines, size, *, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None, font="Calibri"):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = color
    return tb


def add_panel(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_header_bar(slide, x, y, w, h, fill_rgb, text, size=38, color=DARK):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return sh


def add_image_fit(slide, x, y, max_w, max_h, path):
    iw, ih = Image.open(path).size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h, w = max_h, max_h * ratio
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    return slide.shapes.add_picture(path, E(cx), E(cy), E(w), E(h))


def _shape_center_inch(sh):
    return (
        (sh.left + sh.width / 2) / EMU_PER_INCH,
        (sh.top + sh.height / 2) / EMU_PER_INCH,
    )


# --- layout ----------------------------------------------------------------
SW, SH = 48.0, 36.0
M = 0.40
GAP = 0.40
HDR_H = 1.20
PAD_X = 0.40

ROW2_Y = 16.70
ROW2_H = 11.38
ROW3_Y = ROW2_Y + ROW2_H + 0.15
ROW3_H = 6.58
REF_Y = ROW3_Y + ROW3_H + 0.12
REF_H = max(0.62, (SH - M) - REF_Y)

# Asymmetric bottom row: wide F (hero figure), narrower G
F_W = 26.90
G_W = SW - 2 * M - GAP - F_W
X_F = M
X_G = M + F_W + GAP

KEEP = {
    "Picture 10", "TextBox 2", "TextBox 3", "TextBox 1",
    "Title 1", "Subtitle 2",
}

BOTTOM_BAND_Y_MIN = ROW3_Y - 0.08


def strip_bottom_content(slide):
    """Remove shapes whose center lies at/under row 3 (keeps A–E)."""
    victims = []
    for sh in list(slide.shapes):
        if sh.name in KEEP:
            continue
        _, cy = _shape_center_inch(sh)
        if cy >= BOTTOM_BAND_Y_MIN:
            victims.append(sh)
    for sh in victims:
        remove_shape(sh)


if not __import__("os").path.isfile(BASE_PPTX):
    raise FileNotFoundError(
        f"{BASE_PPTX} not found. Build or copy your edited poster there first."
    )

shutil.copy(BASE_PPTX, DST)
prs = Presentation(DST)
slide = prs.slides[0]
strip_bottom_content(slide)

# =============================================================================
# F.  WHAT DRIVES CHURN  — large hero image + compact table / notes
# =============================================================================
add_panel(slide, X_F, ROW3_Y, F_W, ROW3_H, F_BG)
add_header_bar(slide, X_F, ROW3_Y, F_W, HDR_H, F_HDR,
               "F.  What Drives Churn", 36)

cx = X_F + PAD_X
cy = ROW3_Y + HDR_H + 0.12
cw = F_W - 2 * PAD_X

fi_max_h = ROW3_Y + ROW3_H - cy - 0.22
fi_w = cw * 0.70
gap_mid = 0.22
add_image_fit(slide, cx, cy, fi_w, fi_max_h,
              f"{FIG}/feature_importance_catboost.png")

ft_x = cx + fi_w + gap_mid
ft_w = cw - fi_w - gap_mid
feat_rows = [
    ["Top driver", "Importance"],
    ["tenure", "21.4"],
    ["Contract  =  Two year", "14.4"],
    ["InternetService  =  Fiber optic", "13.0"],
    ["TotalCharges", " 7.6"],
    ["MonthlyCharges", " 7.4"],
    ["Contract  =  One year", " 4.8"],
    ["PaymentMethod  =  Electronic check", " 4.3"],
]
ft_h = min(fi_max_h * 0.58, 3.15)
ft = slide.shapes.add_table(len(feat_rows), 2,
                            E(ft_x), E(cy), E(ft_w), E(ft_h)).table
for r in range(len(feat_rows)):
    for c in range(2):
        cell = ft.cell(r, c)
        cell.text = feat_rows[r][c]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.name = "Calibri"
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TBL_HDR_BG
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TBL_ALT_BG if r % 2 == 0 else TBL_RW_BG

interp_y = cy + ft_h + 0.08
interp_h = max(0.55, fi_max_h - ft_h - 0.12)
add_text(slide, ft_x, interp_y, ft_w, interp_h, [
    "Aligns with industry intuition:  short tenure, no long contract,",
    "fiber, electronic check  --  known retention friction points.",
    "Lever:  contract upsell campaigns target the top drivers.",
], size=12, color=DARK, italic=True, line_spacing=1.18)

# =============================================================================
# G.  CONCLUSIONS & NEXT STEPS  (narrower column)
# =============================================================================
add_panel(slide, X_G, ROW3_Y, G_W, ROW3_H, G_BG)
add_header_bar(slide, X_G, ROW3_Y, G_W, HDR_H, G_HDR,
               "G.  Conclusions & Next Steps", 32)

cx = X_G + PAD_X
cy = ROW3_Y + HDR_H + 0.12
cw = G_W - 2 * PAD_X

gw_l = cw * 0.56
gw_r = cw - gw_l - 0.18
gx_r = cx + gw_l + 0.18

add_text(slide, cx, cy, gw_l, 0.40, ["Findings"],
         size=18, bold=True, color=NAVY)
add_text(slide, cx, cy + 0.42, gw_l, 4.0, [
    "1.  CatBoost is the strongest classifier (F1 0.624, ROC-AUC 0.841, lowest Brier 0.139).",
    "2.  Custom MLP does not beat logistic regression -- nonlinearity adds little on small one-hot data.",
    "3.  Cost-aware threshold tuning is the dominant lever -- all four models recover ~$370K per 1,000 customers; architecture changes the result by < 1%.",
    "4.  Counterfactuals turn probabilities into specific retention actions.",
], size=11.5, color=DARK, line_spacing=1.22)

add_text(slide, gx_r, cy, gw_r, 0.36, ["Deployment readiness"],
         size=13, bold=True, color=NAVY)
add_text(slide, gx_r, cy + 0.36, gw_r, 1.05, [
    "Save:  model + scaler + threshold (0.31).",
    "Score monthly;  recalibrate quarterly.",
], size=11, color=DARK, line_spacing=1.18)

add_text(slide, gx_r, cy + 1.48, gw_r, 0.36, ["Limitations"],
         size=13, bold=True, color=NAVY)
add_text(slide, gx_r, cy + 1.84, gw_r, 1.05, [
    "Static snapshot;  no temporal drift.",
    "Cost ratios illustrative -- vary by segment.",
], size=11, color=DARK, line_spacing=1.18)

add_text(slide, gx_r, cy + 2.98, gw_r, 0.36, ["Next steps"],
         size=13, bold=True, color=NAVY)
add_text(slide, gx_r, cy + 3.34, gw_r, 1.05, [
    "Stacking ensemble.",
    "Cost-sensitive loss in MLP.",
    "Live A/B test of threshold.",
], size=11, color=DARK, line_spacing=1.18)

# =============================================================================
# REFERENCES STRIP
# =============================================================================
add_panel(slide, M, REF_Y, SW - 2 * M, REF_H, LGRAY, line=GRAY)
add_text(slide, M + 0.28, REF_Y + 0.04, SW - 2 * M - 0.56, 0.32,
         ["References & Acknowledgements"], size=15, bold=True, color=NAVY)
add_text(slide, M + 0.28, REF_Y + 0.36, SW - 2 * M - 0.56, REF_H - 0.40, [
    "[1] Verbeke, W. et al. (2012). New insights into churn prediction in telecommunications. Eur. J. Op. Res., 218(1), 211-229.",
    "[2] Prokhorenkova, L. et al. (2018). CatBoost: unbiased boosting with categorical features. NeurIPS.",
    "[3] Chen, T. & Guestrin, C. (2016). XGBoost: a scalable tree-boosting system. KDD.",
    "Dataset:  IBM Telco Customer Churn (Kaggle: blastchar/telco-customer-churn).    Acknowledgements:  Dr. Edgar Lobaton (instructor, ECE 525, NC State University).",
], size=11, color=DARK, line_spacing=1.08)

prs.save(DST)
print("Wrote", DST, "| bottom band replaced; rows 1-2 unchanged from", BASE_PPTX)
